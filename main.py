import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import xlrd
import sys
import os
import win32com.client


def meritCertis(certificate_name, spreadsheetDirectory, participantNumber):
    # Insert the directory of your spreadsheet (pls use the template inside this directory)
    workbook = xlrd.open_workbook(rf'{spreadsheetDirectory}')

    # Choose the index of your sheet
    sheet = workbook.sheet_by_index(0)

    lst = []
    posDict = {1: '1st', 2: '2nd', 3: '3rd'}

    for i in range(1, participantNumber+2):
        position = sheet.cell_value(i, 3)
        if type(position) != float:
            if position.lower() == 'participation':
                continue
        else:
            pass
            position = int(position)
            position = posDict[position]
        name = sheet.cell(i, 0).value
        lst.append(name)
        a = name.split()
        firstName, lastName = a[0].capitalize(), a[1].capitalize()
        event = sheet.cell(i, 2).value

        pres = Presentation()
        layout = pres.slide_layouts[6]

        slide = pres.slides.add_slide(layout)
        pic = slide.shapes.add_picture(
            f'{certificate_name}', left=0, top=0, height=pres.slide_height, width=pres.slide_width)
        nameTextBox = slide.shapes.add_textbox(left=Inches(
            0.66), top=Inches(2.23), width=Inches(1), height=Inches(1))
        nameTF = nameTextBox.text_frame
        p1 = nameTF.add_paragraph()
        p1.text = f'{firstName}\n{lastName}'
        p1.font.size = Pt(66)
        p1.font.name = 'Airbnb Cereal App'
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(20, 202, 98)

        positionTextBox = slide.shapes.add_textbox(left=Inches(
            1.815), top=Inches(5.11), width=Inches(1), height=Inches(1))
        posTF = positionTextBox.text_frame
        p2 = posTF.add_paragraph()
        p2.text = position
        p2.font.size = Pt(14)
        p2.font.name = 'Airbnb Cereal App'
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0, 0, 0)

        eventTextBox = slide.shapes.add_textbox(left=Inches(
            4.1), top=Inches(5.11), width=Inches(1), height=Inches(1))
        eventTF = eventTextBox.text_frame
        p3 = eventTF.add_paragraph()
        p3.text = f" {event}"
        p3.font.size = Pt(14)
        p3.font.name = 'Airbnb Cereal App'
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(0, 0, 0)

        if lst.count(name) > 1:
            fileName = f'merit-{name}-{lst.count()}.pptx'
        else:
            fileName = f'merit-{name}.pptx'
        pres.save(fileName)
        print(name, event, position, sep='\n')


def convertMeritToPDF():
    # %% Get path of running script
    script_path = sys.argv[0]

    # %% Get real path
    real_path = os.path.realpath(script_path)

    # %% Get directory path
    folder_path = os.path.dirname(real_path)

    # %% Add final slash at end
    folder_path += "\\"

    # %% Get files in input folder
    input_file_paths = os.listdir(folder_path)

    # %% Convert each file
    for input_file_name in input_file_paths:
        try:
            i = input_file_name.split('-')
        except:
            i = input_file_name

        # Skip files that are not merit-certi ppts
        if i[0] != 'merit' and not input_file_name.lower().endswith((".ppt", ".pptx")):
            continue
        elif i[0] == 'merit' and not input_file_name.lower().endswith((".ppt", ".pptx")):
            continue
        elif i[0] != 'merit' and input_file_name.lower().endswith((".ppt", ".pptx")):
            continue

        # Create input file path
        input_file_path = os.path.join(folder_path, input_file_name)

        # Create powerpoint application object
        powerpoint = win32com.client.DispatchEx("Powerpoint.Application")

        # Set visibility to minimize
        powerpoint.Visible = 1

        # Open the powerpoint slides
        slides = powerpoint.Presentations.Open(input_file_path)

        # Get base file name
        file_name = os.path.splitext(input_file_name)[0]

        # Create output file path
        output_file_path = os.path.join(folder_path, file_name + ".pdf")

        # Save as PDF (formatType = 32)
        slides.SaveAs(output_file_path, 32)

        # Close the slide deck
        slides.Close()


while True:
    a = int(input('\nWhat do you want to do now?'
                  '\n1 -> merit certificates generate'
                  '\n2 -> convert merit certis PPT to PDF'
                  '\n3 -> exit the program'
                  '\n >>> '))

    if a == 1:
        meritCertis(input('enter certificate directory: '), input(
            'enter spreadsheet directory: '), int(input('enter the number of participants: ')))
        print('\nmerit certis ppt generated:)')

    elif a == 2:
        convertMeritToPDF()
        print('\nconverted merit ppts to pdf bb')

    elif a == 3:
        print("Thanks for using certigen. Hope you had a great experience!")
        break
# Add requirements.txt
